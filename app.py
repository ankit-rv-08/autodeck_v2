import streamlit as st
from pptx import Presentation
from pptx.util import Inches
import PyPDF2
import io
import requests

# -----------------------------
# Generate Slide Content via Ollama
# -----------------------------
def generate_slide_content(text, model="llama3"):
    url = "http://localhost:11434/api/chat"
    headers = {"Content-Type": "application/json"}
    payload = {
        "model": model,
        "messages": [
            {"role": "user", "content": f"Summarize this for a slide:\n\n{text}\n\nFormat: Title + 3-5 bullet points."}
        ],
        "stream": False
    }

    try:
        response = requests.post(url, headers=headers, json=payload, timeout=60)
        data = response.json()
        return data.get("message", {}).get("content", "⚠️ LLaMA returned no content.")
    except Exception as e:
        return f"⚠️ Error: {str(e)}"

# -----------------------------
# Extract text from uploaded PDF
# -----------------------------
def extract_text_from_pdf(pdf_file):
    reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# -----------------------------
# Create PowerPoint Presentation
# -----------------------------
def create_presentation(slide_contents, output_filename="autodeck_output.pptx"):
    prs = Presentation()
    for slide_text in slide_contents:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        lines = slide_text.strip().split("\n")
        title = lines[0]
        bullets = lines[1:]
        slide.shapes.title.text = title
        content = slide.placeholders[1]
        for bullet in bullets:
            content.text += f"\n• {bullet.strip()}"
    prs.save(output_filename)
    return output_filename

# -----------------------------
# Streamlit Interface
# -----------------------------
def main():
    st.title("🧠 Autodeck AI v2 — Powered by LLaMA via Ollama")
    st.write("Upload a PDF or paste raw data to auto-generate presentation slides.")

    slide_contents = []

    input_method = st.radio("Select input method:", ("Upload PDF", "Paste Text"))

    text_data = ""
    if input_method == "Upload PDF":
        uploaded_file = st.file_uploader("Upload your PDF", type=["pdf"])
        if uploaded_file:
            text_data = extract_text_from_pdf(uploaded_file)
            st.success("✅ PDF processed successfully!")
    else:
        text_data = st.text_area("Paste raw content below:", height=300)

    if st.button("Generate Slides with AI"):
        if not text_data.strip():
            st.warning("Please upload a PDF or paste some text.")
            return

        st.info("Generating slide content using LLaMA model...")
        chunks = [text_data[i:i+1500] for i in range(0, len(text_data), 1500)]
        for i, chunk in enumerate(chunks):
            with st.spinner(f"Processing chunk {i+1}/{len(chunks)}..."):
                slide_text = generate_slide_content(chunk)
                slide_contents.append(slide_text)

        pptx_file = create_presentation(slide_contents)
        with open(pptx_file, "rb") as f:
            st.download_button(
                label="📥 Download Slides",
                data=f,
                file_name="autodeck_slides.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            st.success("✅ Presentation ready!")

if __name__ == "__main__":
    main()
